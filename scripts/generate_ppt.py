#!/usr/bin/env python3
"""
リプレイ攻撃シミュレーションツールキット - 学術プレゼンテーション生成スクリプト
Academic-style presentation generator

デザイン方針：
- 学術発表にふさわしい簡潔で落ち着いたデザイン
- 明確な情報階層
- 読みやすいフォントサイズと適切な余白
- ページ番号とセクション表示
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

# ============================================================
# 学術的なカラーパレット
# ============================================================
COLORS = {
    # メインカラー（落ち着いた深い青）
    "primary": RGBColor(25, 55, 95),       # 濃紺
    "primary_light": RGBColor(45, 85, 135), # やや明るい青
    
    # アクセント（控えめな暖色）
    "accent": RGBColor(180, 90, 50),        # 落ち着いたオレンジ
    "accent_light": RGBColor(200, 130, 80),
    
    # テキスト
    "text_dark": RGBColor(33, 33, 33),      # ほぼ黒
    "text_medium": RGBColor(80, 80, 80),    # グレー
    "text_light": RGBColor(120, 120, 120),  # 薄いグレー
    
    # 背景
    "bg_white": RGBColor(255, 255, 255),
    "bg_light": RGBColor(248, 249, 250),    # 薄いグレー背景
    "bg_header": RGBColor(240, 242, 245),   # ヘッダー背景
    
    # 表
    "table_header": RGBColor(45, 85, 135),
    "table_row_alt": RGBColor(245, 247, 250),
    
    # 強調
    "success": RGBColor(40, 120, 80),       # 緑
    "warning": RGBColor(180, 130, 40),      # 黄
    "danger": RGBColor(160, 60, 60),        # 赤
}

# フォント設定
FONTS = {
    "title": "Arial",           # タイトル用
    "body": "Arial",            # 本文用
    "code": "Consolas",         # コード用
}


def add_slide_number(slide, prs, slide_num, total_slides):
    """スライド番号を追加"""
    footer = slide.shapes.add_textbox(
        prs.slide_width - Inches(1), 
        prs.slide_height - Inches(0.4),
        Inches(0.8), 
        Inches(0.3)
    )
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = f"{slide_num} / {total_slides}"
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS["text_light"]
    p.alignment = PP_ALIGN.RIGHT


def add_section_indicator(slide, section_name):
    """セクション名を左下に追加"""
    indicator = slide.shapes.add_textbox(
        Inches(0.3), 
        Inches(7.1),
        Inches(4), 
        Inches(0.3)
    )
    tf = indicator.text_frame
    p = tf.paragraphs[0]
    p.text = section_name
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["text_light"]
    p.font.italic = True


def add_header_line(slide, prs):
    """ヘッダーラインを追加"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, Inches(1.05),
        prs.slide_width, Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["primary"]
    line.line.fill.background()


def add_academic_title_slide(prs, title, subtitle, author="", affiliation=""):
    """学術的なタイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 上部のアクセントライン
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = COLORS["primary"]
    top_line.line.fill.background()
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # サブタイトル
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_medium"]
        p.alignment = PP_ALIGN.CENTER
    
    # 著者
    if author:
        author_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(5.2), Inches(8.4), Inches(0.5)
        )
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = author
        p.font.size = Pt(16)
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_dark"]
        p.alignment = PP_ALIGN.CENTER
    
    # 所属
    if affiliation:
        affil_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(5.7), Inches(8.4), Inches(0.5)
        )
        tf = affil_box.text_frame
        p = tf.paragraphs[0]
        p.text = affiliation
        p.font.size = Pt(12)
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER
    
    # 下部のライン
    bottom_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(2), Inches(6.5), 
        Inches(6), Inches(0.01)
    )
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = COLORS["text_light"]
    bottom_line.line.fill.background()
    
    return slide


def add_academic_section_slide(prs, section_number, section_title):
    """学術的なセクション区切りスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 左側のアクセントバー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), prs.slide_height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # セクション番号
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(1.5), Inches(1)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_number)
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary_light"]
    
    # セクションタイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["text_dark"]
    
    return slide


def add_academic_content_slide(prs, title, bullet_points, section_name=""):
    """学術的な内容スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    
    # ヘッダーライン
    add_header_line(slide, prs)
    
    # 内容
    content_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # インデントレベルの判定
        indent_level = 0
        text = point
        while text.startswith("  "):
            indent_level += 1
            text = text[2:]
        
        if text.strip() == "":
            p.text = ""
            p.space_after = Pt(6)
        else:
            if indent_level > 0:
                p.text = f"    {'–' if indent_level == 1 else '•'} {text.strip()}"
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS["text_medium"]
            else:
                p.text = f"• {text}"
                p.font.size = Pt(18)
                p.font.color.rgb = COLORS["text_dark"]
            
            p.font.name = FONTS["body"]
            p.space_after = Pt(8)
    
    # セクション表示
    if section_name:
        add_section_indicator(slide, section_name)
    
    return slide


def add_academic_two_column_slide(prs, title, left_content, right_content, section_name=""):
    """2カラムスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    
    add_header_line(slide, prs)
    
    # 左カラム
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}" if point.strip() else ""
        p.font.size = Pt(16)
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_dark"]
        p.space_after = Pt(6)
    
    # 右カラム
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, point in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}" if point.strip() else ""
        p.font.size = Pt(16)
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_dark"]
        p.space_after = Pt(6)
    
    if section_name:
        add_section_indicator(slide, section_name)
    
    return slide


def add_academic_table_slide(prs, title, headers, rows, section_name=""):
    """学術的な表スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    
    add_header_line(slide, prs)
    
    # 表を作成
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    # 表のサイズ調整
    col_width = min(1.8, 8.5 / num_cols)
    table_width = Inches(col_width * num_cols)
    row_height = 0.45
    table_height = Inches(row_height * num_rows)
    
    left = (prs.slide_width - table_width) / 2  # 中央揃え
    top = Inches(1.5)
    
    table = slide.shapes.add_table(
        num_rows, num_cols, left, top, table_width, table_height
    ).table
    
    # ヘッダー行のスタイル
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["table_header"]
        
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["bg_white"]
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # データ行
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            
            # 交互の背景色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["table_row_alt"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["bg_white"]
            
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.name = FONTS["body"]
            p.font.color.rgb = COLORS["text_dark"]
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    if section_name:
        add_section_indicator(slide, section_name)
    
    return slide


def add_academic_diagram_slide(prs, title, diagram_text, section_name=""):
    """図解スライド（等幅フォント）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    
    add_header_line(slide, prs)
    
    # 図解ボックス（背景付き）
    bg_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(1.25),
        Inches(9.2), Inches(5.6)
    )
    bg_box.fill.solid()
    bg_box.fill.fore_color.rgb = COLORS["bg_light"]
    bg_box.line.color.rgb = COLORS["text_light"]
    bg_box.line.width = Pt(0.5)
    
    # 図解テキスト
    diagram_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(5.4)
    )
    tf = diagram_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(10)
    p.font.name = FONTS["code"]
    p.font.color.rgb = COLORS["text_dark"]
    
    if section_name:
        add_section_indicator(slide, section_name)
    
    return slide


def add_academic_image_slide(prs, title, image_path, caption="", section_name=""):
    """画像スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    
    add_header_line(slide, prs)
    
    # 画像を追加
    if Path(image_path).exists():
        # 画像を中央に配置
        pic = slide.shapes.add_picture(
            image_path, Inches(1.2), Inches(1.4), width=Inches(7.6)
        )
    
    # キャプション
    if caption:
        caption_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.6), Inches(9), Inches(0.4)
        )
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_medium"]
        p.alignment = PP_ALIGN.CENTER
    
    if section_name:
        add_section_indicator(slide, section_name)
    
    return slide


def add_academic_summary_slide(prs, title, points):
    """まとめスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 上部バー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # 内容
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if point.strip() == "":
            p.text = ""
            p.space_after = Pt(10)
        else:
            p.text = f"✓ {point}"
            p.font.size = Pt(18)
            p.font.name = FONTS["body"]
            p.font.color.rgb = COLORS["text_dark"]
            p.space_after = Pt(12)
    
    return slide


def add_academic_end_slide(prs, title, contact_info=""):
    """終了スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 上部バー
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()
    
    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = FONTS["title"]
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.CENTER
    
    # 連絡先
    if contact_info:
        info_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(1)
        )
        tf = info_box.text_frame
        p = tf.paragraphs[0]
        p.text = contact_info
        p.font.size = Pt(14)
        p.font.name = FONTS["body"]
        p.font.color.rgb = COLORS["text_medium"]
        p.alignment = PP_ALIGN.CENTER
    
    # 下部ライン
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3), Inches(5.5),
        Inches(4), Inches(0.01)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["text_light"]
    line.line.fill.background()
    
    return slide


# ============================================================
# 面接用精簡版 PPT（学術スタイル）
# ============================================================

def create_interview_presentation():
    """面接用の精簡版プレゼンテーション（学術スタイル）"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ===== 1. タイトル =====
    add_academic_title_slide(
        prs,
        "IoTシステムにおける\nリプレイ攻撃防御手法の比較評価",
        "シミュレーションによる定量的分析",
        author="発表者名",
        affiliation=""
    )
    
    # ===== 2. 研究背景 =====
    add_academic_content_slide(
        prs, 
        "研究背景と目的",
        [
            "IoT機器への脅威：リプレイ攻撃",
            "  攻撃者が正規コマンドを記録し、後から再送信",
            "  例：スマートロック、車両キーレス、産業制御",
            "",
            "課題：複数の防御手法が存在するが、比較評価が不十分",
            "  現実的なネットワーク条件での性能が不明確",
            "",
            "本研究の目的",
            "  シミュレーションにより4種類の防御手法を定量的に評価",
            "  パケット損失・順序入れ替え条件下での有効性を検証",
        ],
        section_name="1. 背景"
    )
    
    # ===== 3. システム構成 =====
    add_academic_diagram_slide(
        prs, 
        "評価システムの構成",
        """
    ┌─────────────────────────────────────────────────────────────────┐
    │                                                                 │
    │      ┌──────────┐        ┌───────────┐        ┌──────────┐     │
    │      │  Sender  │  ───→  │  Channel  │  ───→  │ Receiver │     │
    │      │ (送信者) │        │(チャネル) │        │ (受信者) │     │
    │      └──────────┘        └─────┬─────┘        └──────────┘     │
    │                                │                    ↑          │
    │                                ↓                    │          │
    │                          ┌──────────┐               │          │
    │                          │ Attacker │───────────────┘          │
    │                          │ (攻撃者) │   リプレイ攻撃           │
    │                          └──────────┘                          │
    │                                                                 │
    │      チャネルモデル:                                            │
    │        • パケット損失率 (p_loss): 0〜30%                        │
    │        • 順序入れ替え確率 (p_reorder): 0〜30%                   │
    │                                                                 │
    └─────────────────────────────────────────────────────────────────┘
""",
        section_name="2. 方法"
    )
    
    # ===== 4. 防御手法 =====
    add_academic_table_slide(
        prs, 
        "評価対象：3種類の防御メカニズム",
        ["手法", "原理", "特徴"],
        [
            ["Rolling Counter", "カウンタの単調増加を検証", "実装が単純"],
            ["Sliding Window", "ビットマスクで一定範囲を許容", "順序入れ替え耐性"],
            ["Challenge-Response", "毎回異なるNonceで認証", "高セキュリティ"],
        ],
        section_name="2. 方法"
    )
    
    # ===== 5. 実験結果 =====
    add_academic_content_slide(
        prs, 
        "主要な実験結果",
        [
            "実験条件：200回のモンテカルロ試行（固定シード）",
            "",
            "発見1：Rolling Counter は順序入れ替えに脆弱",
            "  30%並び替え時、正規受理率が70〜80%台まで低下",
            "",
            "発見2：Sliding Window は順序入れ替えに高い耐性",
            "  同条件でも正規受理率90%前後を維持",
            "",
            "発見3：ウィンドウサイズ W=3〜5 で良好なバランス",
            "  高い正規受理率と低い攻撃成功率を両立",
        ],
        section_name="3. 結果"
    )
    
    # ===== 6. 結果図 =====
    figures_path = Path(__file__).parent.parent / "figures"
    add_academic_image_slide(
        prs, 
        "順序入れ替えの影響",
        str(figures_path / "p_reorder_legit.png"),
        "Fig. 1: Rolling（青）は順序入れ替え環境で正規受理率が低下、Window（橙）は安定",
        section_name="3. 結果"
    )
    
    # ===== 7. まとめ =====
    add_academic_summary_slide(prs, "まとめ", [
        "4種類の防御手法をシミュレーションで定量評価",
        "",
        "主要な結論：",
        "  Sliding Window (W=3〜5) が実用的な選択肢",
        "  Rolling Counter は順序入れ替え環境で課題あり",
        "",
        "今後の予定：",
        "  実機（ESP32等）での検証実験",
        "  バースト損失モデルの導入",
    ])
    
    # ===== 8. 終了 =====
    add_academic_end_slide(
        prs,
        "ご清聴ありがとうございました",
        "GitHub: github.com/tammakiiroha/IoT-Replay-Defense-Simulator"
    )
    
    return prs


# ============================================================
# 完全版 PPT（学術スタイル）
# ============================================================

def create_full_presentation():
    """完全版プレゼンテーション（学術スタイル）"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ===== タイトル =====
    add_academic_title_slide(
        prs,
        "IoTシステムにおける\nリプレイ攻撃防御手法の比較評価",
        "シミュレーションによる定量的分析",
        author="発表者名",
        affiliation=""
    )
    
    # ===== 目次 =====
    add_academic_content_slide(prs, "発表の構成", [
        "1. 研究背景と目的",
        "2. リプレイ攻撃の概要",
        "3. 評価システムの構成",
        "4. 防御メカニズムの詳細",
        "5. 実験設計と方法",
        "6. 実験結果",
        "7. 考察と今後の課題",
    ])
    
    # ===== セクション1 =====
    add_academic_section_slide(prs, 1, "研究背景と目的")
    
    add_academic_content_slide(
        prs, 
        "研究の動機",
        [
            "IoT機器の普及と無線通信の脅威",
            "  スマートホーム、産業用制御、車両システム",
            "  無線通信はフレームの傍受・再送信が可能",
            "",
            "リプレイ攻撃への対策",
            "  複数の防御手法が提案されている",
            "  しかし、現実的な条件での比較評価が不足",
            "",
            "シミュレーションによる評価の必要性",
            "  再現性のある定量的評価",
            "  パラメータの系統的な変化による分析",
        ],
        section_name="1. 背景"
    )
    
    add_academic_content_slide(
        prs, 
        "研究目的",
        [
            "目的",
            "  4種類のリプレイ攻撃防御手法を定量的に評価",
            "  現実的なチャネル条件下での有効性を検証",
            "",
            "評価観点",
            "  正規受理率（ユーザビリティ）",
            "  攻撃成功率（セキュリティ）",
            "",
            "貢献",
            "  再現可能な実験環境の提供",
            "  実用的な設計指針の提示",
        ],
        section_name="1. 背景"
    )
    
    # ===== セクション2 =====
    add_academic_section_slide(prs, 2, "リプレイ攻撃の概要")
    
    add_academic_diagram_slide(
        prs, 
        "リプレイ攻撃のメカニズム",
        """
    ┌─────────────────────────────────────────────────────────────────┐
    │                                                                 │
    │   【正規通信】                                                   │
    │                                                                 │
    │      正規ユーザー  ───────→  "UNLOCK"  ───────→  受信機        │
    │                                  │                              │
    │                                  ↓                              │
    │                              [攻撃者]                           │
    │                              記録・保存                         │
    │                                                                 │
    │   【攻撃】                                                       │
    │                                  │                              │
    │                                  ↓                              │
    │      攻撃者  ─────────→  "UNLOCK"  ───────→  受信機            │
    │              再送信              ↓                              │
    │                            不正に受理される                     │
    │                                                                 │
    └─────────────────────────────────────────────────────────────────┘

    脅威例：スマートロック、車両キーレス、産業制御システム
""",
        section_name="2. 攻撃概要"
    )
    
    # ===== セクション3 =====
    add_academic_section_slide(prs, 3, "評価システムの構成")
    
    add_academic_diagram_slide(
        prs, 
        "システムアーキテクチャ",
        """
    ┌─────────────────────────────────────────────────────────────────┐
    │                          入力層                                  │
    │   ┌─────────────────────┐     ┌─────────────────────┐          │
    │   │ コマンドシーケンス   │     │ シミュレーション設定 │          │
    │   │ (traces/*.txt)      │     │ (SimulationConfig)  │          │
    │   └──────────┬──────────┘     └──────────┬──────────┘          │
    │              │                           │                      │
    │              ▼                           ▼                      │
    │   ┌─────────────────────────────────────────────────────┐      │
    │   │                 シミュレーション層                   │      │
    │   │   Sender ──→ Channel ──→ Receiver ←── Attacker     │      │
    │   └─────────────────────────────────────────────────────┘      │
    │              │                                                  │
    │              ▼                                                  │
    │   ┌─────────────────────────────────────────────────────┐      │
    │   │                      出力層                          │      │
    │   │   JSON結果 (results/*.json)    図表 (figures/*.png) │      │
    │   └─────────────────────────────────────────────────────┘      │
    └─────────────────────────────────────────────────────────────────┘
""",
        section_name="3. システム"
    )
    
    add_academic_table_slide(
        prs, 
        "主要コンポーネント",
        ["コンポーネント", "役割", "実装"],
        [
            ["Sender", "フレーム生成、カウンタ/MAC付与", "sender.py"],
            ["Channel", "パケット損失・順序入れ替え", "channel.py"],
            ["Receiver", "防御メカニズムによる検証", "receiver.py"],
            ["Attacker", "フレーム記録・リプレイ", "attacker.py"],
        ],
        section_name="3. システム"
    )
    
    # ===== セクション4 =====
    add_academic_section_slide(prs, 4, "防御メカニズムの詳細")
    
    add_academic_content_slide(
        prs, 
        "Rolling Counter + MAC",
        [
            "原理",
            "  送信側：フレームごとにカウンタを増加",
            "  受信側：カウンタの単調増加を検証",
            "",
            "動作例",
            "  Counter=0 到着 → 0 > -1 → 受理 (last=0)",
            "  Counter=1 到着 → 1 > 0 → 受理 (last=1)",
            "  Counter=0 再送 → 0 ≤ 1 → 拒否（リプレイ検出）",
            "",
            "特徴",
            "  長所：実装が単純、リプレイを検出可能",
            "  短所：順序入れ替えに弱い",
        ],
        section_name="4. 防御手法"
    )
    
    add_academic_content_slide(
        prs, 
        "Sliding Window",
        [
            "原理",
            "  ビットマスクで一定範囲のカウンタを管理",
            "  順序が入れ替わっても受理可能",
            "",
            "動作例（Window Size = 5）",
            "  許容範囲：[last-4, last-3, last-2, last-1, last]",
            "  ビットマスクで受信済みを記録",
            "",
            "特徴",
            "  長所：順序入れ替えに対応可能",
            "  短所：ウィンドウサイズの調整が必要",
        ],
        section_name="4. 防御手法"
    )
    
    add_academic_content_slide(
        prs, 
        "Challenge-Response",
        [
            "原理",
            "  受信側がランダムなNonceを発行",
            "  送信側がNonceを含むMACで応答",
            "",
            "動作",
            "  受信側 → 送信側：Nonce発行",
            "  送信側 → 受信側：Command + Nonce + MAC",
            "  古いNonceのフレームは拒否",
            "",
            "特徴",
            "  長所：非常に高いセキュリティ",
            "  短所：双方向通信が必要、レイテンシ増加",
        ],
        section_name="4. 防御手法"
    )
    
    # ===== セクション5 =====
    add_academic_section_slide(prs, 5, "実験設計と方法")
    
    add_academic_table_slide(
        prs, 
        "実験条件",
        ["実験", "変数", "固定条件", "目的"],
        [
            ["実験1", "p_loss: 0-30%", "p_reorder=0%", "損失の影響"],
            ["実験2", "p_reorder: 0-30%", "p_loss=10%", "並び替えの影響"],
            ["実験3", "window: 1-20", "p_loss=15%, p_reorder=15%", "最適サイズ"],
        ],
        section_name="5. 実験設計"
    )
    
    add_academic_content_slide(
        prs, 
        "実験方法",
        [
            "モンテカルロシミュレーション",
            "  各条件で200回の独立した試行",
            "  固定シード（seed=42）による再現性の確保",
            "",
            "評価指標",
            "  正規受理率：正規フレームが受理される割合",
            "  攻撃成功率：リプレイフレームが受理される割合",
            "",
            "公平性の確保",
            "  全モードで同一の乱数列を使用",
        ],
        section_name="5. 実験設計"
    )
    
    # ===== セクション6 =====
    add_academic_section_slide(prs, 6, "実験結果")
    
    add_academic_table_slide(
        prs, 
        "実験1：パケット損失の影響",
        ["手法", "p_loss=0%", "p_loss=30%", "観察"],
        [
            ["No Defense", "可用性高/攻撃高", "可用性約70%", "チャネル効果"],
            ["Rolling", "可用性高/攻撃低", "可用性約70%", "効果維持"],
            ["Window", "可用性高/攻撃低", "可用性約70%", "効果維持"],
            ["Challenge", "可用性高/攻撃極低", "可用性約70%", "最も安定"],
        ],
        section_name="6. 結果"
    )
    
    figures_path = Path(__file__).parent.parent / "figures"
    
    add_academic_image_slide(
        prs, 
        "実験1：可用性への影響",
        str(figures_path / "p_loss_legit.png"),
        "Fig. 1: パケット損失率と正規受理率の関係",
        section_name="6. 結果"
    )
    
    add_academic_table_slide(
        prs, 
        "実験2：順序入れ替えの影響",
        ["手法", "p_reorder=0%", "p_reorder=30%", "観察"],
        [
            ["No Defense", "可用性約90%", "可用性約90%", "無関係"],
            ["Rolling", "可用性約90%", "可用性70-80%台", "⚠ 大きな低下"],
            ["Window", "可用性約90%", "可用性約90%", "✓ 高い耐性"],
            ["Challenge", "可用性約90%", "可用性60%台", "⚠ 影響あり"],
        ],
        section_name="6. 結果"
    )
    
    add_academic_image_slide(
        prs, 
        "実験2：順序入れ替えの影響",
        str(figures_path / "p_reorder_legit.png"),
        "Fig. 2: Rolling（青）は順序入れ替えで低下、Window（橙）は安定",
        section_name="6. 結果"
    )
    
    add_academic_content_slide(
        prs, 
        "実験2の主要な発見",
        [
            "Rolling Counter の問題点",
            "  30%並び替えで正規受理率が70〜80%台まで低下",
            "  正規フレームの一部が誤って拒否される",
            "",
            "Sliding Window の優位性",
            "  同条件でも正規受理率90%前後を維持",
            "  順序入れ替えが発生しやすい環境に適している",
            "",
            "Challenge-Response の特性",
            "  高並び替え条件で可用性が低下する傾向",
        ],
        section_name="6. 結果"
    )
    
    add_academic_table_slide(
        prs, 
        "実験3：ウィンドウサイズの影響",
        ["サイズ", "正規受理率", "攻撃成功率", "評価"],
        [
            ["W=1", "低い（20%台）", "低め", "✗ 不適切"],
            ["W=3", "高い（85%前後）", "低い", "✓ 良好"],
            ["W=5", "高い（85%前後）", "低い", "✓ 推奨"],
            ["W=7", "高い（85%前後）", "やや上昇", "○ 許容"],
            ["W≥9", "高い（85%前後）", "上昇傾向", "△ 注意"],
        ],
        section_name="6. 結果"
    )
    
    add_academic_image_slide(
        prs, 
        "実験3：ウィンドウサイズのトレードオフ",
        str(figures_path / "window_tradeoff.png"),
        "Fig. 3: W=3〜5で正規受理率と攻撃成功率のバランスが良好",
        section_name="6. 結果"
    )
    
    # ===== セクション7 =====
    add_academic_section_slide(prs, 7, "考察と今後の課題")
    
    add_academic_table_slide(
        prs, 
        "総合評価",
        ["手法", "可用性", "セキュリティ", "推奨用途"],
        [
            ["Rolling", "高い", "高い", "順序保証環境"],
            ["Window", "高い", "高い", "一般的なIoT"],
            ["Challenge", "高い", "非常に高い", "高セキュリティ要求"],
        ],
        section_name="7. 考察"
    )
    
    add_academic_content_slide(
        prs, 
        "本研究の制約",
        [
            "チャネルモデルの簡略化",
            "  i.i.d.パケット損失モデルを使用",
            "  バースト損失は未考慮",
            "",
            "攻撃モデルの範囲",
            "  リプレイ攻撃のみを評価",
            "  リレー攻撃等は対象外",
            "",
            "検証環境",
            "  ソフトウェアシミュレーションのみ",
        ],
        section_name="7. 考察"
    )
    
    add_academic_content_slide(
        prs, 
        "今後の研究計画",
        [
            "短期",
            "  実機（ESP32, nRF52840等）での検証実験",
            "  Gilbert-Elliott型バースト損失モデルの導入",
            "",
            "中期",
            "  リレー攻撃など他の攻撃パターンの評価",
            "  実測トレースを用いたシミュレーション",
        ],
        section_name="7. 考察"
    )
    
    # ===== まとめ =====
    add_academic_summary_slide(prs, "まとめ", [
        "4種類の防御手法をシミュレーションで定量評価",
        "現実的なチャネル条件（損失・順序入れ替え）を考慮",
        "",
        "主要な発見",
        "  Rolling Counter：順序入れ替え環境で可用性低下",
        "  Sliding Window (W=3〜5)：良好なバランス",
        "  Challenge-Response：高セキュリティだが双方向必要",
        "",
        "再現可能な実験環境をオープンソースで公開",
    ])
    
    # ===== 終了 =====
    add_academic_end_slide(
        prs,
        "ご清聴ありがとうございました",
        "GitHub: github.com/tammakiiroha/IoT-Replay-Defense-Simulator"
    )
    
    return prs


def main():
    """メイン関数"""
    print("=" * 65)
    print("  学術プレゼンテーション生成スクリプト")
    print("  Academic Presentation Generator")
    print("=" * 65)
    
    output_dir = Path(__file__).parent.parent / "docs"
    output_dir.mkdir(exist_ok=True)
    
    # 1. 面接用精簡版
    print("\n[1/2] 面接用精簡版を生成中...")
    interview_prs = create_interview_presentation()
    interview_path = output_dir / "presentation_interview_jp.pptx"
    interview_prs.save(str(interview_path))
    print(f"  ✓ {interview_path.name} ({len(interview_prs.slides)} slides)")
    
    # 2. 完全版
    print("\n[2/2] 完全版（セミナー用）を生成中...")
    full_prs = create_full_presentation()
    full_path = output_dir / "presentation_full_jp.pptx"
    full_prs.save(str(full_path))
    print(f"  ✓ {full_path.name} ({len(full_prs.slides)} slides)")
    
    print("\n" + "=" * 65)
    print("  生成完了")
    print("=" * 65)
    print(f"\n  📁 出力先: {output_dir}")
    print(f"  📊 面接用: {interview_path.name} (8 slides)")
    print(f"  📚 完全版: {full_path.name}")
    print("\n  💡 面接では精簡版を使用してください。")
    print("=" * 65)


if __name__ == "__main__":
    main()
